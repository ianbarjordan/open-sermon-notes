"""
01_ingest_files.py — Walk source dir, parse every file, apply quarantine pipeline,
write JSON to data/documents/.

Usage:
    python build/01_ingest_files.py --source SampleData/ --dry-run --verbose
    python build/01_ingest_files.py --source SampleData/ --verbose
    python build/01_ingest_files.py --source SampleData/ --limit 10
"""
import argparse
import hashlib
import json
import os
import platform
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from app.config import MIN_CHUNK_WORDS  # noqa: E402

# ---------------------------------------------------------------------------
# Persistent dedup registry
# ---------------------------------------------------------------------------
PROCESSED_REGISTRY = "data/processed.json"


def load_registry(registry_path: str) -> dict:
    """Load {source_path: sha256} registry from disk. Returns empty dict if missing."""
    p = Path(registry_path)
    if p.exists():
        try:
            with open(p, encoding='utf-8') as fh:
                return json.load(fh)
        except Exception as e:
            print(f"WARNING: could not load registry {registry_path}: {e}", file=sys.stderr)
    return {}


def save_registry(registry: dict, registry_path: str) -> None:
    """Atomically write the registry back to disk."""
    p = Path(registry_path)
    p.parent.mkdir(parents=True, exist_ok=True)
    tmp_fd, tmp_path = tempfile.mkstemp(dir=str(p.parent), suffix='.tmp')
    try:
        with os.fdopen(tmp_fd, 'w', encoding='utf-8') as fh:
            json.dump(registry, fh, indent=2, ensure_ascii=False)
        os.replace(tmp_path, str(p))
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise
from build.format_detect import detect_format  # noqa: E402
from build.parse_filename import parse_filename  # noqa: E402

# ---------------------------------------------------------------------------
# Admin / non-faith keyword lists
# ---------------------------------------------------------------------------
ADMIN_KEYWORDS = [
    'tax', 'minute', 'apportionment', 'budget', 'invoice',
    'bank holiday', 'reimburs', 'agenda', 'commitment card',
    'obituary', 'job description', 'giving card', 'philmont', 'scout',
]

FAITH_KEYWORDS = [
    'jesus', 'christ', 'lord', 'god', 'gospel', 'grace', 'scripture',
    'bible', 'sermon', 'prayer', 'faith', 'worship', 'salvation',
    'amen', 'resurrection', 'parable', 'disciple', 'holy spirit',
    'covenant', 'redemption', 'forgiveness', 'eternal', 'heaven',
]

# Extensions to skip silently
SILENT_SKIP_SUFFIXES = {'.identifier', '.csv', '.md'}


class ParseError(Exception):
    pass


# ---------------------------------------------------------------------------
# File collection
# ---------------------------------------------------------------------------

def collect_files(source_dir: str) -> list[Path]:
    """Return sorted list of all files under source_dir."""
    root = Path(source_dir)
    return sorted(p for p in root.rglob('*') if p.is_file())


# ---------------------------------------------------------------------------
# Quarantine filters
# ---------------------------------------------------------------------------

def should_skip_silently(path: Path) -> bool:
    return path.suffix.lower() in SILENT_SKIP_SUFFIXES


def is_filename_flagged(stem: str) -> bool:
    lower = stem.lower()
    return any(kw in lower for kw in ADMIN_KEYWORDS)


# ---------------------------------------------------------------------------
# OLE2 parsers — platform-aware
# ---------------------------------------------------------------------------

# Persistent Word COM instance (Windows only); avoids starting Word per file.
_WIN_WORD_APP = None


def _get_win_word_app():
    """Lazy-initialize and return a persistent Word COM instance (Windows only)."""
    global _WIN_WORD_APP
    if _WIN_WORD_APP is None:
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            raise ParseError(
                "pywin32 is required on Windows to parse .doc files.\n"
                "Install with:  pip install pywin32\n"
                "Also requires Microsoft Word to be installed."
            )
        try:
            pythoncom.CoInitialize()
            _WIN_WORD_APP = win32com.client.Dispatch("Word.Application")
            _WIN_WORD_APP.Visible = False
            _WIN_WORD_APP.DisplayAlerts = 0
        except Exception as e:
            raise ParseError(
                f"Could not start Microsoft Word COM automation: {e}\n"
                "Ensure Microsoft Word is installed."
            )
    return _WIN_WORD_APP


def _close_win_word_app() -> None:
    """Quit the persistent Word COM instance if active."""
    global _WIN_WORD_APP
    if _WIN_WORD_APP is not None:
        try:
            _WIN_WORD_APP.Quit()
        except Exception:
            pass
        _WIN_WORD_APP = None


def _parse_ole2_windows(path: Path) -> str:
    """Windows: extract text from OLE2 .doc via Word COM automation (pywin32)."""
    word = _get_win_word_app()
    abs_path = str(path.resolve())
    doc = None
    try:
        doc = word.Documents.Open(
            abs_path,
            ConfirmConversions=False,
            ReadOnly=True,
            AddToRecentFiles=False,
        )
        text = doc.Content.Text
    except Exception as e:
        msg = str(e)
        if 'detected a problem' in msg or '-2146821993' in msg:
            raise ParseError(
                f"Word blocked {path.name} (security/trust error).\n"
                "Fix: In Word → Options → Trust Center → Trusted Locations → "
                "add your sermon folder. Or run in PowerShell: "
                "Get-ChildItem 'E:\\your\\sermons' -Recurse | Unblock-File"
            ) from e
        raise ParseError(f"Word COM failed on {path.name}: {e}") from e
    finally:
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                pass
    if not text.strip():
        raise ParseError("Word COM produced empty output")
    return text


def _parse_ole2_antiword(path: Path) -> str:
    """Linux/macOS: extract text from OLE2 .doc via antiword."""
    result = subprocess.run(
        ['antiword', '-t', str(path)],
        capture_output=True, text=True, timeout=30,
    )
    if result.returncode not in (0, 1):  # antiword returns 1 on minor warnings
        raise ParseError(f"antiword failed (rc={result.returncode}): {result.stderr[:200]}")
    text = result.stdout
    if not text.strip():
        raise ParseError("antiword produced empty output")
    return text


def parse_ole2(path: Path) -> str:
    """Extract plain text from an OLE2 Word file.

    Dispatches to Word COM automation on Windows or antiword on Linux/macOS.
    """
    if platform.system() == 'Windows':
        return _parse_ole2_windows(path)
    return _parse_ole2_antiword(path)


# ---------------------------------------------------------------------------
# Other parsers
# ---------------------------------------------------------------------------

def _strip_rtf(raw: str) -> str:
    """Minimal RTF stripper: no oletools needed.

    RTF files use cp1252 encoding (Windows standard for Word 97-2003 RTF),
    which is why we decode bytes as cp1252 before calling this function.
    """
    # Remove \* escaped groups (annotations, stylesheet, etc.)
    text = re.sub(r'\\\*[^{}]*', '', raw)
    # Remove { } groups that are purely control (non-text content blocks)
    text = re.sub(r'\{[^{}]{0,200}\}', '', text)
    # Decode \'XX hex escapes (cp1252)
    def decode_hex(m: re.Match) -> str:
        try:
            return bytes.fromhex(m.group(1)).decode('cp1252', errors='replace')
        except ValueError:
            return ''
    text = re.sub(r"\\'([0-9a-fA-F]{2})", decode_hex, text)
    # Remove control words \word and control symbols \X
    text = re.sub(r'\\[a-zA-Z]+\d*\s?', '', text)
    text = re.sub(r'\\[^a-zA-Z]', '', text)
    # Strip remaining braces
    text = text.replace('{', '').replace('}', '')
    # Normalize whitespace and line endings (handles \r\n, \r, \n)
    text = re.sub(r'\r\n|\r|\n', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def parse_rtf(path: Path) -> str:
    # RTF files from Word 97-2003 are always cp1252 encoded
    raw = path.read_bytes().decode('cp1252', errors='replace')
    text = _strip_rtf(raw)
    if not text.strip():
        raise ParseError("RTF stripper produced empty output")
    return text


def parse_docx(path: Path) -> str:
    try:
        import docx  # python-docx
    except ImportError:
        raise ParseError("python-docx not installed")
    try:
        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = '\n\n'.join(paragraphs)
    except Exception as e:
        raise ParseError(f"python-docx failed: {e}") from e
    if not text.strip():
        raise ParseError("python-docx produced empty output")
    return text


def _get_pptx_presentation(path: Path):
    """Return (prs, text) for a pptx file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ParseError("python-pptx not installed")
    try:
        prs = Presentation(str(path))
        slides_text = []
        for slide in prs.slides:
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = ' '.join(r.text for r in para.runs).strip()
                        if line:
                            parts.append(line)
            if parts:
                slides_text.append('\n'.join(parts))
    except Exception as e:
        raise ParseError(f"python-pptx failed: {e}") from e
    return prs, '\n\n'.join(slides_text)


# ---------------------------------------------------------------------------
# PPTX slide heuristics
# ---------------------------------------------------------------------------

def _slide_lines(slide) -> list[str]:
    lines = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = ' '.join(r.text for r in para.runs).strip()
                if line:
                    lines.append(line)
    return lines


def is_worship_slides(prs) -> bool:
    """Return True if >80% of slides are 'short-line only' (worship/song slides)."""
    total = len(prs.slides)
    if total == 0:
        return False
    short_count = 0
    for slide in prs.slides:
        lines = _slide_lines(slide)
        if not lines:
            continue
        if all(len(line.split()) <= 8 for line in lines):
            short_count += 1
    return (short_count / total) > 0.80


def is_sparse_pptx(prs) -> bool:
    """Return True if fewer than 3 slides have any text content."""
    text_slides = sum(
        1 for slide in prs.slides if any(_slide_lines(slide))
    )
    return text_slides < 3


# ---------------------------------------------------------------------------
# Content analysis
# ---------------------------------------------------------------------------

def count_faith_hits(text: str) -> int:
    lower = text.lower()
    return sum(1 for kw in FAITH_KEYWORDS if kw in lower)


def compute_sha256(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, 'rb') as fh:
        for block in iter(lambda: fh.read(65536), b''):
            h.update(block)
    return h.hexdigest()


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def make_doc_id(stem: str) -> str:
    """Sanitize a filename stem into a valid doc_id."""
    doc_id = stem.lower()
    doc_id = re.sub(r'[^a-z0-9_]', '_', doc_id)
    doc_id = re.sub(r'_+', '_', doc_id)
    doc_id = doc_id.strip('_')
    return doc_id


def quarantine(path: Path, reason: str, root: str, dry_run: bool = False) -> None:
    """Copy (not move) file to quarantine/<reason>/ subdirectory.

    Silently skips the copy if the destination already exists or is not writable
    (e.g. a re-run over a partially-completed ingest).
    """
    dest_dir = Path(root) / reason
    if not dry_run:
        try:
            dest_dir.mkdir(parents=True, exist_ok=True)
            dest = dest_dir / path.name
            if dest.exists():
                return  # already quarantined from a previous run
            shutil.copy2(str(path), str(dest))
        except PermissionError as e:
            print(f"  WARNING: could not quarantine {path.name}: {e}", file=sys.stderr)


def write_document_json(doc: dict, out_dir: str) -> None:
    """Atomically write document JSON via a temp file."""
    out_path = Path(out_dir) / f"{doc['doc_id']}.json"
    tmp_fd, tmp_path = tempfile.mkstemp(dir=out_dir, suffix='.tmp')
    try:
        with os.fdopen(tmp_fd, 'w', encoding='utf-8') as fh:
            json.dump(doc, fh, indent=2, ensure_ascii=False)
        os.replace(tmp_path, out_path)
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise


# ---------------------------------------------------------------------------
# Core ingest function
# ---------------------------------------------------------------------------

def ingest_file(
    path: Path,
    out_dir: str,
    quarantine_root: str,
    seen_hashes: set,
    dry_run: bool = False,
    force: bool = False,
    verbose: bool = False,
    registry: dict | None = None,
) -> str:
    """Ingest a single file through the quarantine pipeline.

    Returns an outcome string (one of the quarantine reasons or 'accepted').
    """
    stem = path.stem

    # Step 1: Silent skip
    if should_skip_silently(path):
        if verbose:
            print(f"  SKIP    {path.name}")
        return 'skipped'

    # Step 2: .pub files
    if path.suffix.lower() == '.pub':
        if verbose:
            print(f"  PUB     {path.name}")
        quarantine(path, 'format_pub', quarantine_root, dry_run)
        return 'format_pub'

    # Step 3: Admin filename keyword
    if is_filename_flagged(stem):
        if verbose:
            print(f"  ADMIN   {path.name}")
        quarantine(path, 'filename_flagged', quarantine_root, dry_run)
        return 'filename_flagged'

    # Determine format
    fmt = detect_format(str(path))
    ext = path.suffix.lower()

    # Step 4: Parse dispatch
    try:
        if fmt == 'ooxml' and ext == '.docx':
            text = parse_docx(path)
            doc_format = 'ooxml_doc'
            prs = None
        elif fmt == 'ooxml' and ext == '.pptx':
            prs, text = _get_pptx_presentation(path)
            doc_format = 'ooxml_pptx'
        elif fmt == 'ole2':
            text = parse_ole2(path)
            doc_format = 'ole2'
            prs = None
        elif fmt == 'rtf':
            text = parse_rtf(path)
            doc_format = 'rtf'
            prs = None
        else:
            raise ParseError(f"Unknown format: {fmt} (ext={ext})")
    except ParseError as e:
        if verbose:
            print(f"  PARSE_ERR {path.name}: {e}")
        quarantine(path, 'manual_review', quarantine_root, dry_run)
        return 'manual_review'

    # Step 5: Word count check
    wc = len(text.split())
    if wc < MIN_CHUNK_WORDS:
        if verbose:
            print(f"  SHORT   {path.name} ({wc} words)")
        quarantine(path, 'too_short', quarantine_root, dry_run)
        return 'too_short'

    # Steps 6-7: PPTX slide heuristics
    if doc_format == 'ooxml_pptx' and prs is not None:
        if is_worship_slides(prs):
            if verbose:
                print(f"  WORSHIP {path.name}")
            quarantine(path, 'worship_slides', quarantine_root, dry_run)
            return 'worship_slides'
        if is_sparse_pptx(prs):
            if verbose:
                print(f"  SPARSE  {path.name}")
            quarantine(path, 'sparse_pptx', quarantine_root, dry_run)
            return 'sparse_pptx'

    # Step 8: Faith content check
    hits = count_faith_hits(text)
    if hits < 2:
        if verbose:
            print(f"  NONFAITH {path.name} ({hits} hits)")
        quarantine(path, 'non_faith', quarantine_root, dry_run)
        return 'non_faith'

    # Step 9: Duplicate check
    sha = compute_sha256(path)
    if sha in seen_hashes:
        if verbose:
            print(f"  DUP     {path.name}")
        quarantine(path, 'duplicates', quarantine_root, dry_run)
        return 'duplicates'
    seen_hashes.add(sha)
    if registry is not None and not dry_run:
        registry[path.as_posix()] = sha

    # Step 10: Accepted — build document JSON
    meta = parse_filename(stem)
    doc_id = make_doc_id(stem)

    # Store as POSIX path (forward slashes) for cross-platform consistency
    source_file = path.as_posix()

    doc = {
        'doc_id':        doc_id,
        'source_file':   source_file,
        'title':         meta.get('title'),
        'scripture_ref': meta.get('scripture_ref'),
        'date':          meta.get('date'),
        'format':        doc_format,
        'word_count':    wc,
        'text':          text,
    }

    if not dry_run:
        Path(out_dir).mkdir(parents=True, exist_ok=True)
        out_path = Path(out_dir) / f"{doc_id}.json"
        if out_path.exists() and not force:
            if verbose:
                print(f"  SKIP(exists) {path.name}")
            return 'skipped_exists'
        write_document_json(doc, out_dir)

    if verbose:
        print(f"  OK      {path.name} ({wc} words)")
    return 'accepted'


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(
        description='Ingest sermon files through quarantine pipeline.'
    )
    parser.add_argument('--source',     metavar='PATH', default='SampleData',
                        help='Source directory (default: SampleData)')
    parser.add_argument('--out',        metavar='PATH', default='data/documents',
                        help='JSON output dir (default: data/documents)')
    parser.add_argument('--quarantine', metavar='PATH', default='raw/quarantine',
                        help='Quarantine root (default: raw/quarantine)')
    parser.add_argument('--limit',      metavar='INT',  type=int, default=0,
                        help='Stop after N files (0 = no limit)')
    parser.add_argument('--dry-run',    action='store_true',
                        help='Full pipeline, no files written')
    parser.add_argument('--force',      action='store_true',
                        help='Re-process files that already have JSON output')
    parser.add_argument('--verbose',    action='store_true',
                        help='Per-file decisions to stdout')
    parser.add_argument('--registry',   metavar='PATH', default=PROCESSED_REGISTRY,
                        help=f'Persistent hash registry (default: {PROCESSED_REGISTRY})')
    args = parser.parse_args()

    files = collect_files(args.source)
    if args.limit:
        files = files[:args.limit]

    print(f"Found {len(files)} files in {args.source!r}")
    if platform.system() == 'Windows':
        print("Platform: Windows — .doc files will use Word COM automation (pywin32)")
    else:
        print(f"Platform: {platform.system()} — .doc files will use antiword")
    if args.dry_run:
        print("DRY RUN — no files will be written")

    # Load persistent registry and pre-populate seen_hashes
    registry: dict = load_registry(args.registry) if not args.dry_run else {}
    seen_hashes: set = set(registry.values())
    print(f"Registry: {len(registry)} previously accepted files")

    counts: dict[str, int] = {}

    try:
        from tqdm import tqdm
        iterator = tqdm(files, unit='file')
    except ImportError:
        iterator = files  # type: ignore[assignment]

    try:
        for path in iterator:
            outcome = ingest_file(
                path=path,
                out_dir=args.out,
                quarantine_root=args.quarantine,
                seen_hashes=seen_hashes,
                dry_run=args.dry_run,
                force=args.force,
                verbose=args.verbose,
                registry=registry,
            )
            counts[outcome] = counts.get(outcome, 0) + 1
    finally:
        # Always close the Word COM instance if it was opened on Windows
        _close_win_word_app()

    # Save updated registry
    if not args.dry_run:
        save_registry(registry, args.registry)
        print(f"Registry saved: {len(registry)} entries → {args.registry}")

    print("\n--- Ingest Summary ---")
    total = sum(counts.values())
    for outcome, n in sorted(counts.items(), key=lambda x: -x[1]):
        print(f"  {outcome:25s} {n:4d}")
    print(f"  {'TOTAL':25s} {total:4d}")


if __name__ == '__main__':
    main()
